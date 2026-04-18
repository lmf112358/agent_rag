"""
文档处理模块
支持多种文档格式加载和智能分块策略
"""

import os
import re
import logging
from typing import List, Optional, Dict, Any, Callable,Tuple
from pathlib import Path
from langchain_core.documents import Document

logger = logging.getLogger(__name__)


class BaseLoader:
    """基础文档加载器"""

    def load(self) -> List[Document]:
        raise NotImplementedError


class TextLoader(BaseLoader):
    """文本加载器（内置，无第三方依赖）"""

    def __init__(self, file_path: str, encoding: str = "utf-8", *args, **kwargs):
        self.file_path = file_path
        self.encoding = encoding

    def load(self) -> List[Document]:
        with open(self.file_path, "r", encoding=self.encoding, errors="ignore") as file:
            content = file.read()
        return [Document(page_content=content, metadata={"source": self.file_path})]


class CSVLoader(BaseLoader):
    """CSV加载器（内置，无第三方依赖）"""

    def __init__(self, file_path: str, encoding: str = "utf-8", *args, **kwargs):
        self.file_path = file_path
        self.encoding = encoding

    def load(self) -> List[Document]:
        import csv

        rows: List[str] = []
        with open(self.file_path, "r", encoding=self.encoding, errors="ignore", newline="") as file:
            reader = csv.reader(file)
            for row in reader:
                rows.append(",".join(row))

        return [Document(page_content="\n".join(rows), metadata={"source": self.file_path})]


class PyPDFLoader(BaseLoader):
    """PDF 加载器（使用 pypdf）"""

    def __init__(self, file_path: str, *args, **kwargs):
        self.file_path = file_path

    def load(self) -> List[Document]:
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError("pypdf is required. Install it: pip install pypdf")

        reader = PdfReader(self.file_path)
        docs = []
        for page_num, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                metadata = {
                    "source": self.file_path,
                    "page": page_num + 1,
                }
                docs.append(Document(page_content=text, metadata=metadata))
        return docs


class UnstructuredWordDocumentLoader(BaseLoader):
    """DOCX 加载器（使用 python-docx）"""

    def __init__(self, file_path: str, *args, **kwargs):
        self.file_path = file_path

    def load(self) -> List[Document]:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError("python-docx is required. Install it: pip install python-docx")

        doc = DocxDocument(self.file_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)

        text = "\n".join(full_text)
        if not text.strip():
            return []

        return [Document(page_content=text, metadata={"source": self.file_path})]


class DocLoader(BaseLoader):
    """旧版 DOC 加载器（支持多种方式）"""

    def __init__(self, file_path: str, *args, **kwargs):
        self.file_path = file_path

    def load(self) -> List[Document]:
        """尝试多种方式读取 .doc 文件"""
        text = None
        methods_tried = []

        # 方法 1: 尝试使用 textract（如果已安装）
        try:
            import textract
            methods_tried.append("textract")
            text = textract.process(self.file_path).decode("utf-8", errors="ignore")
        except ImportError:
            pass
        except Exception as e:
            logger.debug(f"textract failed: {e}")

        # 方法 2: 尝试使用 antiword（命令行工具）
        if text is None:
            try:
                import subprocess
                methods_tried.append("antiword")
                result = subprocess.run(
                    ["antiword", self.file_path],
                    capture_output=True,
                    text=True,
                    timeout=30,
                )
                if result.returncode == 0:
                    text = result.stdout
            except (ImportError, FileNotFoundError, subprocess.TimeoutExpired):
                pass
            except Exception as e:
                logger.debug(f"antiword failed: {e}")

        # 方法 3: Windows 平台尝试使用 pywin32
        if text is None:
            try:
                import sys
                if sys.platform == "win32":
                    methods_tried.append("pywin32")
                    import win32com.client
                    word = win32com.client.Dispatch("Word.Application")
                    word.Visible = False
                    try:
                        doc = word.Documents.Open(os.path.abspath(self.file_path))
                        text = doc.Content.Text
                        doc.Close()
                    finally:
                        word.Quit()
            except ImportError:
                pass
            except Exception as e:
                logger.debug(f"pywin32 failed: {e}")

        # 方法 4: 尝试使用 docx2txt
        if text is None:
            try:
                import docx2txt
                methods_tried.append("docx2txt")
                text = docx2txt.process(self.file_path)
            except ImportError:
                pass
            except Exception as e:
                logger.debug(f"docx2txt failed: {e}")

        if text is None:
            raise ImportError(
                f"Failed to read .doc file. Tried methods: {methods_tried}. "
                f"Please install one of: textract, antiword, pywin32 (Windows), docx2txt, "
                f"or convert the file to .docx format manually."
            )

        if not text.strip():
            return []

        return [Document(page_content=text, metadata={"source": self.file_path})]


class UnstructuredExcelLoader(BaseLoader):
    """Excel 加载器（可选，提示安装）"""

    def __init__(self, file_path: str, *args, **kwargs):
        self.file_path = file_path

    def load(self) -> List[Document]:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("pandas is required for Excel. Install it: pip install pandas")

        df = pd.read_excel(self.file_path)
        text = df.to_csv(sep="\t", na_rep="")
        return [Document(page_content=text, metadata={"source": self.file_path})]


class UnstructuredPowerPointLoader(BaseLoader):
    """PPT 加载器（可选，提示安装）"""

    def __init__(self, file_path: str, *args, **kwargs):
        self.file_path = file_path

    def load(self) -> List[Document]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PowerPoint. Install it: pip install python-pptx")

        prs = Presentation(self.file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)

        text = "\n".join(texts)
        if not text.strip():
            return []

        return [Document(page_content=text, metadata={"source": self.file_path})]


class RecursiveCharacterTextSplitter:
    """简化递归文本分割器（无 langchain-text-splitters 依赖）"""

    def __init__(
        self,
        chunk_size: int = 512,
        chunk_overlap: int = 100,
        separators: Optional[List[str]] = None,
        length_function: Callable[[str], int] = len,
        add_start_index: bool = False,
        **kwargs,
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = separators or ["\n\n", "\n", " ", ""]
        self.length_function = length_function
        self.add_start_index = add_start_index

    def split_text(self, text: str) -> List[str]:
        if not text:
            return []

        chunks: List[str] = []
        start = 0
        total_len = len(text)

        while start < total_len:
            max_end = min(start + self.chunk_size, total_len)
            end = max_end

            if max_end < total_len:
                window = text[start:max_end]
                for sep in self.separators:
                    if not sep:
                        continue
                    idx = window.rfind(sep)
                    if idx > 0:
                        end = start + idx + len(sep)
                        break

            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            if end >= total_len:
                break

            start = max(end - self.chunk_overlap, start + 1)

        return chunks

    def split_documents(self, documents: List[Document]) -> List[Document]:
        split_docs: List[Document] = []

        for doc in documents:
            chunks = self.split_text(doc.page_content)
            cursor = 0

            for chunk in chunks:
                metadata = dict(doc.metadata) if doc.metadata else {}
                if self.add_start_index:
                    idx = doc.page_content.find(chunk, cursor)
                    if idx < 0:
                        idx = cursor
                    metadata["start_index"] = idx
                    cursor = idx + len(chunk)

                split_docs.append(Document(page_content=chunk, metadata=metadata))

        return split_docs


class CharacterTextSplitter(RecursiveCharacterTextSplitter):
    """兼容占位：与递归分割器保持一致"""


class ChunkConfig:
    """分块配置"""

    def __init__(
        self,
        chunk_size: int = 512,
        chunk_overlap: int = 100,
        separators: Optional[List[str]] = None,
        length_function: Callable[[str], int] = len,
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = separators or ["\n\n", "\n", "。", "；", " ", ""]
        self.length_function = length_function


class DocumentProcessor:
    """文档处理器"""

    LOADER_MAPPING: Dict[str, type] = {
        ".pdf": PyPDFLoader,
        ".docx": UnstructuredWordDocumentLoader,
        ".doc": DocLoader,
        ".xlsx": UnstructuredExcelLoader,
        ".xls": UnstructuredExcelLoader,
        ".pptx": UnstructuredPowerPointLoader,
        ".ppt": UnstructuredPowerPointLoader,
        ".csv": CSVLoader,
        ".txt": TextLoader,
        ".md": TextLoader,
    }

    def __init__(
        self,
        chunk_config: Optional[ChunkConfig] = None,
        extract_images: bool = False,
        use_mineru: bool = True,
    ):
        self.chunk_config = chunk_config or ChunkConfig()
        self.extract_images = extract_images
        self.use_mineru = use_mineru
        self.mineru_client = None
        if self.use_mineru:
            self._init_mineru()

    def _init_mineru(self):
        """初始化 MinerU 客户端"""
        try:
            from langchain_rag.document.mineru_client import create_mineru_client_from_config
            self.mineru_client = create_mineru_client_from_config()
            if self.mineru_client:
                logger.info("MinerU client initialized")
            else:
                logger.info("MinerU not enabled in config, using PyPDF for PDF")
        except Exception as e:
            logger.warning(f"Failed to initialize MinerU client: {e}, using PyPDF for PDF")
            self.mineru_client = None

    def get_loader(self, file_path: str) -> BaseLoader:
        """根据文件扩展名获取合适的加载器"""
        ext = Path(file_path).suffix.lower()

        if ext not in self.LOADER_MAPPING:
            raise ValueError(
                f"Unsupported file type: {ext}. "
                f"Supported types: {list(self.LOADER_MAPPING.keys())}"
            )

        loader_class = self.LOADER_MAPPING[ext]

        if ext == ".docx":
            return loader_class(file_path, mode="elements")
        elif ext == ".doc":
            return loader_class(file_path)
        elif ext == ".pdf":
            return loader_class(file_path)
        elif ext in [".xlsx", ".xls"]:
            return loader_class(file_path)
        elif ext == ".csv":
            return loader_class(file_path, encoding="utf-8")
        else:
            return loader_class(file_path)

    def load_document(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> List[Document]:
        """加载单个文档（PDF 优先用 MinerU，失败回退 PyPDF）"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        path = Path(file_path)
        ext = path.suffix.lower()

        # PDF 优先尝试 MinerU
        if ext == ".pdf" and self.use_mineru and self.mineru_client:
            try:
                docs = self._try_load_with_mineru(file_path, metadata)
                if docs:
                    logger.info(f"Loaded PDF with MinerU: {path.name}")
                    return docs
            except Exception as e:
                logger.warning(f"MinerU failed for {path.name}: {e}, falling back to PyPDF")

        # 回退到默认加载器
        loader = self.get_loader(file_path)
        docs = loader.load()

        if metadata:
            for doc in docs:
                doc.metadata.update(metadata)

        return docs

    def _try_load_with_mineru(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> Optional[List[Document]]:
        """尝试用 MinerU 加载 PDF"""
        try:
            from langchain_rag.document.mineru_loader import MinerULoader
        except ImportError:
            logger.debug("MinerULoader not available")
            return None

        try:
            loader = MinerULoader(file_path, client=self.mineru_client)
            docs = loader.load()

            if docs and metadata:
                for doc in docs:
                    doc.metadata.update(metadata)

            return docs
        except Exception as e:
            logger.debug(f"MinerULoader failed: {e}")
            return None

    def load_documents(
        self,
        directory: str,
        glob_pattern: str = "**/*",
        metadata_fn: Optional[Callable[[str], Dict[str, Any]]] = None,
    ) -> List[Document]:
        """批量加载目录下的文档"""
        docs = []
        path = Path(directory)

        for file_path in path.glob(glob_pattern):
            if file_path.is_file():
                try:
                    file_metadata = {}
                    if metadata_fn:
                        file_metadata = metadata_fn(str(file_path))
                    docs.extend(self.load_document(str(file_path), file_metadata))
                except Exception as e:
                    print(f"Warning: Failed to load {file_path}: {str(e)}")

        return docs

    def split_documents(
        self,
        documents: List[Document],
        chunk_config: Optional[ChunkConfig] = None,
    ) -> List[Document]:
        """分割文档为块（表格感知）"""
        cfg = chunk_config or self.chunk_config
        split_docs: List[Document] = []

        for doc in documents:
            content = doc.page_content
            metadata = dict(doc.metadata) if doc.metadata else {}

            # 检测是否包含 Markdown 表格
            if "|" in content and "\n|" in content:
                # 尝试表格感知分块
                table_chunks = self._split_table_aware(content, metadata, cfg)
                split_docs.extend(table_chunks)
            else:
                # 普通文本分块
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=cfg.chunk_size,
                    chunk_overlap=cfg.chunk_overlap,
                    separators=cfg.separators,
                    length_function=cfg.length_function,
                    add_start_index=True,
                )
                chunks = text_splitter.split_text(content)
                for i, chunk in enumerate(chunks):
                    chunk_meta = dict(metadata)
                    chunk_meta["chunk_type"] = "text"
                    chunk_meta["chunk_index"] = i
                    split_docs.append(Document(page_content=chunk, metadata=chunk_meta))

        return split_docs

    def _split_table_aware(
        self,
        content: str,
        metadata: Dict[str, Any],
        cfg: ChunkConfig,
    ) -> List[Document]:
        """表格感知分块：整表保留 + 文本语义分块"""
        chunks: List[Document] = []

        # 简单实现：先按 Markdown 表格边界分割
        lines = content.split("\n")
        current_text = []
        in_table = False
        current_table = []

        for line in lines:
            # 检测表格行（以 | 开头或包含 |）
            is_table_line = line.strip().startswith("|") or (line.count("|") >= 2)

            if is_table_line:
                if not in_table:
                    # 先输出之前的文本
                    if current_text:
                        self._add_text_chunk(current_text, metadata, chunks, cfg)
                        current_text = []
                    in_table = True
                current_table.append(line)
            else:
                if in_table:
                    # 输出表格
                    self._add_table_chunk(current_table, metadata, chunks)
                    current_table = []
                    in_table = False
                current_text.append(line)

        # 输出剩余内容
        if in_table and current_table:
            self._add_table_chunk(current_table, metadata, chunks)
        elif current_text:
            self._add_text_chunk(current_text, metadata, chunks, cfg)

        return chunks

    def _add_text_chunk(
        self,
        lines: List[str],
        base_meta: Dict[str, Any],
        chunks: List[Document],
        cfg: ChunkConfig,
    ):
        """添加文本块"""
        text = "\n".join(lines).strip()
        if not text:
            return

        # 文本太长时分块
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=cfg.chunk_size,
            chunk_overlap=cfg.chunk_overlap,
            separators=cfg.separators,
            length_function=cfg.length_function,
            add_start_index=True,
        )
        text_chunks = text_splitter.split_text(text)
        for i, chunk in enumerate(text_chunks):
            meta = dict(base_meta)
            meta["chunk_type"] = "text"
            meta["chunk_index"] = len(chunks)
            chunks.append(Document(page_content=chunk, metadata=meta))

    def _add_table_chunk(
        self,
        lines: List[str],
        base_meta: Dict[str, Any],
        chunks: List[Document],
    ):
        """添加表格块（整表保留）"""
        table_content = "\n".join(lines).strip()
        if not table_content:
            return

        # 超大表降级：按字符数截断（保留表头）
        MAX_TABLE_CHARS = 2000
        if len(table_content) > MAX_TABLE_CHARS:
            # 尝试找到表头
            header_lines = []
            body_lines = []
            for i, line in enumerate(lines):
                if i < 3:
                    header_lines.append(line)
                else:
                    body_lines.append(line)

            # 每 10 行一个 chunk，重复表头
            rows_per_chunk = 10
            for i in range(0, len(body_lines), rows_per_chunk):
                chunk_lines = header_lines + body_lines[i:i+rows_per_chunk]
                chunk_content = "\n".join(chunk_lines)
                meta = dict(base_meta)
                meta["chunk_type"] = "parameter_table"
                meta["chunk_index"] = len(chunks)
                chunks.append(Document(page_content=chunk_content, metadata=meta))
        else:
            # 整表保留
            meta = dict(base_meta)
            meta["chunk_type"] = "parameter_table"
            meta["chunk_index"] = len(chunks)
            chunks.append(Document(page_content=table_content, metadata=meta))


class ChineseTextSplitter(RecursiveCharacterTextSplitter):
    """针对中文优化的文本分割器"""

    def __init__(
        self,
        chunk_size: int = 512,
        chunk_overlap: int = 100,
        separators: Optional[List[str]] = None,
    ):
        if separators is None:
            separators = [
                "\n\n", "\n", "。", "；", "！", "？", "，", " ",
                "、", "）", "」", "'", "\"", "】", "]]",
            ]
        super().__init__(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=separators,
            length_function=len,
            add_start_index=True,
        )

    def split_text(self, text: str) -> List[str]:
        """分割文本"""
        # 委托给父类的递归分割逻辑，仅在需要时对超长块做中文友好二次切割
        parent_chunks = super().split_text(text)
        result: List[str] = []
        for chunk in parent_chunks:
            if len(chunk) <= self.chunk_size:
                result.append(chunk)
            else:
                result.extend(self._safe_split_to_list(chunk))
        return result

    def _safe_split_to_list(self, text: str) -> List[str]:
        """将超长文本按中文标点安全切割为若干块"""
        if len(text) <= self.chunk_size:
            return [text]

        chunks = []
        start = 0
        while start < len(text):
            end = start + self.chunk_size
            if end < len(text):
                for sep in ["。", "；", "，", "、"]:
                    idx = text.rfind(sep, start, end)
                    if idx > start:
                        end = idx + 1
                        break
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            start = end
        return chunks if chunks else [text]


class MarkdownProcessor:
    """Markdown文档专用处理器"""

    def __init__(self, chunk_size: int = 512, chunk_overlap: int = 100):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def process(self, content: str, metadata: Optional[Dict[str, Any]] = None) -> List[Document]:
        """处理Markdown内容"""
        sections = self._split_by_headers(content)

        docs = []
        for section in sections:
            heading = section.get("heading", "")
            content_text = section.get("content", "").strip()

            if not content_text:
                continue

            chunk_size = self.chunk_size
            content_text = content_text.replace("\r\n", "\n").replace("\r", "\n")

            while len(content_text) > chunk_size:
                chunk = content_text[:chunk_size]
                last_break = max(
                    chunk.rfind("\n"),
                    chunk.rfind("。"),
                    chunk.rfind("；"),
                )
                if last_break > chunk_size // 2:
                    chunk = chunk[: last_break + 1]

                chunk_metadata = {"heading": heading} if heading else {}
                if metadata:
                    chunk_metadata.update(metadata)
                chunk_metadata["chunk_index"] = len(docs)

                docs.append(Document(page_content=chunk, metadata=chunk_metadata))
                content_text = content_text[len(chunk):]
            else:
                if content_text:
                    chunk_metadata = {"heading": heading} if heading else {}
                    if metadata:
                        chunk_metadata.update(metadata)

                    docs.append(Document(page_content=content_text, metadata=chunk_metadata))

        return docs

    def _split_by_headers(self, content: str) -> List[Dict[str, str]]:
        """按标题分割Markdown"""
        lines = content.split("\n")
        sections = []
        current_heading = ""
        current_content = []

        for line in lines:
            header_match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if header_match:
                if current_content:
                    sections.append({
                        "heading": current_heading,
                        "content": "\n".join(current_content),
                    })
                    current_content = []
                current_heading = header_match.group(2)
            else:
                current_content.append(line)

        if current_content:
            sections.append({
                "heading": current_heading,
                "content": "\n".join(current_content),
            })

        return sections


class DocumentMetadata:
    """文档元数据工具类"""

    @staticmethod
    def from_file_path(file_path: str) -> Dict[str, Any]:
        """从文件路径提取元数据"""
        path = Path(file_path)
        return {
            "source": str(path.absolute()),
            "filename": path.name,
            "file_extension": path.suffix,
            "file_size": os.path.getsize(file_path),
        }

    @staticmethod
    def from_path_advanced(file_path: str, root_dir: str = "data") -> Dict[str, Any]:
        """
        从路径提取高级元数据（基于珠海深联真实目录结构）

        路径示例: data/珠海深联高效机房资料20241024/EQP-设备技术资料/EQP-01 冷水机组/特灵---10-22/CCTV - CCTV-1650RT-6.45 - Product Report.pdf
        """
        path = Path(file_path)
        root = Path(root_dir)

        # 基础元数据
        metadata = DocumentMetadata.from_file_path(file_path)

        # 计算相对路径部分
        try:
            rel_parts = list(path.relative_to(root).parts)
        except ValueError:
            rel_parts = []

        # L1: project_name (去掉日期后缀)
        if len(rel_parts) >= 1:
            project_dir = rel_parts[0]
            # 去掉日期后缀: "珠海深联高效机房资料20241024" -> "珠海深联高效机房"
            project_name = re.sub(r'资料\d{8}$', '', project_dir)
            project_name = re.sub(r'\d{8}$', '', project_name)
            metadata["project_name"] = project_name

        # L2: doc_type (去掉前缀: "EQP-设备技术资料" -> "设备技术资料")
        if len(rel_parts) >= 2:
            doc_type_dir = rel_parts[1]
            doc_type = re.sub(r'^[A-Z]+-', '', doc_type_dir)
            metadata["doc_type"] = doc_type

        # L3: equipment_category (去掉前缀: "EQP-01 冷水机组" -> "冷水机组")
        if len(rel_parts) >= 3:
            eq_cat_dir = rel_parts[2]
            eq_cat = re.sub(r'^[A-Z]+-\d+\s*', '', eq_cat_dir)
            metadata["equipment_category"] = eq_cat

        # L4: brand (如果存在第四级目录)
        if len(rel_parts) >= 4:
            brand_dir = rel_parts[3]
            brand = DocumentMetadata._extract_brand(brand_dir)
            if brand:
                metadata["brand"] = brand

        # 从文件名提取 model_spec
        model_spec, model_spec_raw = DocumentMetadata._extract_model_spec(path.name)
        if model_spec:
            metadata["model_spec"] = model_spec
            metadata["model_spec_raw"] = model_spec_raw

        # 从文件名提取 file_type_tag
        file_type_tag = DocumentMetadata._extract_file_type_tag(path.name)
        if file_type_tag:
            metadata["file_type_tag"] = file_type_tag

        return metadata

    @staticmethod
    def _extract_brand(brand_dir: str) -> Optional[str]:
        """从目录名提取品牌"""
        BRAND_PATTERNS = {
            "特灵": ["特灵", "Trane", "CCTV"],
            "开利": ["开利", "Carrier", "19XR"],
            "约克": ["约克", "York"],
            "麦克维尔": ["麦克维尔", "McQuay"],
            "良机": ["良机", "LiangChi", "LCP"],
            "元亨": ["元亨", "Yuanheng"],
            "凯泉": ["凯泉", "Kaiquan"],
        }
        for brand, patterns in BRAND_PATTERNS.items():
            for p in patterns:
                if p in brand_dir:
                    return brand
        return None

    @staticmethod
    def _extract_model_spec(filename: str) -> Tuple[Optional[str], str]:
        """从文件名提取型号（标准化 + 原始）"""
        raw = filename
        # 尝试正则模式
        MODEL_PATTERNS = [
            # 特灵模式: CCTV-1650RT-6.45
            r'(CCTV)[-\s]?(\d+)RT[-\s]?(\d+\.\d+)',
            # 开利模式: 19XR-84V4F30MHT5A
            r'(19XR)[-\s]?([A-Z0-9]+)',
            # 良机模式: LCP-4059S-L-C1-JC
            r'(LCP)[-\s]?([A-Z0-9-]+)',
            # SRN 模式: SRN-900LG-1
            r'(SRN)[-\s]?([A-Z0-9-]+)',
            # 通用模式: 字母-数字-字母
            r'([A-Z]{2,})[-\s]?([A-Z0-9-]+)',
        ]
        for pattern in MODEL_PATTERNS:
            match = re.search(pattern, filename, re.IGNORECASE)
            if match:
                groups = match.groups()
                normalized = "-".join(groups).upper()
                return normalized, raw
        # fallback: 清理空格和特殊字符
        cleaned = re.sub(r'[^\w\-]', '', filename)
        return None, raw

    @staticmethod
    def _extract_file_type_tag(filename: str) -> Optional[str]:
        """从文件名提取文件类型标签"""
        TAGS = ["参数表", "样本", "Product Report", "外形图", "基础图", "报价", "IPLV", "选型报告"]
        for tag in TAGS:
            if tag in filename:
                return tag
        return None

    @staticmethod
    def add_document_type(
        metadata: Dict[str, Any],
        doc_type: str,
    ) -> Dict[str, Any]:
        """添加文档类型"""
        metadata["document_type"] = doc_type
        return metadata

    @staticmethod
    def add_classification(
        metadata: Dict[str, Any],
        classification: str,
    ) -> Dict[str, Any]:
        """添加分类标签"""
        metadata["classification"] = classification
        return metadata

    @staticmethod
    def add_project_info(
        metadata: Dict[str, Any],
        project_name: str,
        project_year: Optional[str] = None,
    ) -> Dict[str, Any]:
        """添加项目信息"""
        metadata["project_name"] = project_name
        if project_year:
            metadata["project_year"] = project_year
        return metadata


def load_and_process_documents(
    file_paths: List[str],
    chunk_size: int = 512,
    chunk_overlap: int = 100,
    document_type: Optional[str] = None,
    add_metadata: Optional[Dict[str, Any]] = None,
) -> List[Document]:
    """便捷函数：加载并处理文档"""
    processor = DocumentProcessor(
        chunk_config=ChunkConfig(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
    )

    all_docs = []
    for file_path in file_paths:
        metadata = DocumentMetadata.from_file_path(file_path)
        if document_type:
            metadata = DocumentMetadata.add_document_type(metadata, document_type)
        if add_metadata:
            metadata.update(add_metadata)

        docs = processor.load_document(file_path, metadata)
        all_docs.extend(docs)

    return processor.split_documents(all_docs)
